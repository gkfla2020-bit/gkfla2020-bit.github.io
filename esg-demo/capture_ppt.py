import asyncio
from pathlib import Path
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu

SLIDE_DIR = Path('slides')
SLIDE_DIR.mkdir(exist_ok=True)

URL = 'http://localhost:8080'
WIDTH = 1920
HEIGHT = 1080

async def capture_slides():
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(
            viewport={'width': WIDTH, 'height': HEIGHT},
            device_scale_factor=1
        )
        await page.goto(URL, wait_until='networkidle')
        await page.wait_for_timeout(4000)

        count = await page.evaluate('''() => {
            document.querySelectorAll('.ppt-slide').forEach(s => {
                var c = s.firstElementChild;
                if (c && !c.classList.contains('slide-num')) {
                    c.style.transform = 'none';
                }
            });
            // 슬라이드 번호 숨기기
            document.querySelectorAll('.slide-num').forEach(e => e.style.display='none');
            return document.querySelectorAll('.ppt-slide').length;
        }''')
        print(f'Found {count} slides')

        for i in range(count):
            await page.evaluate(f'''() => {{
                var s = document.querySelectorAll('.ppt-slide')[{i}];
                if (s) s.scrollIntoView({{behavior:'instant'}});
            }}''')
            await page.wait_for_timeout(600)
            path = SLIDE_DIR / f'slide_{i+1:02d}.png'
            await page.screenshot(path=str(path), type='png')
            print(f'  Captured {i+1}/{count}')

        await browser.close()
    return count

def make_pptx(count):
    # PPT 슬라이드 크기: 정확히 16:9
    SLIDE_W = Emu(12192000)  # 13.333 inches
    SLIDE_H = Emu(6858000)   # 7.5 inches

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i in range(count):
        img = SLIDE_DIR / f'slide_{i+1:02d}.png'
        if not img.exists():
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 이미지를 슬라이드 전체에 꽉 차게
        slide.shapes.add_picture(
            str(img),
            left=0,
            top=0,
            width=SLIDE_W,
            height=SLIDE_H
        )

    # 동영상 슬라이드 추가
    video_path = Path('demo-video.mp4')
    if video_path.exists():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 동영상도 꽉 차게
        slide.shapes.add_movie(
            str(video_path),
            left=0, top=0,
            width=SLIDE_W, height=SLIDE_H,
            poster_frame_image=str(SLIDE_DIR / 'slide_01.png'),
            mime_type='video/mp4'
        )

    out = 'ESG_TradeGuard_Presentation.pptx'
    prs.save(out)
    print(f'\nDone: {out} ({count} slides + video)')

if __name__ == '__main__':
    count = asyncio.run(capture_slides())
    make_pptx(count)
