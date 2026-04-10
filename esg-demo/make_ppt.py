from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HANA = RGBColor(0x00, 0x8C, 0x73)
HANA_DARK = RGBColor(0x00, 0x4D, 0x3D)
HANA_LIGHT = RGBColor(0xE6, 0xF5, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1B, 0x1B, 0x1B)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0x9C, 0xA3, 0xAF)
BG = RGBColor(0xF7, 0xF8, 0xFA)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name='Noto Sans KR'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_label(slide, left, top, text):
    add_text(slide, left, top, 3, 0.3, text, size=9, color=HANA, bold=True)

def add_title(slide, left, top, text, width=10):
    add_text(slide, left, top, width, 0.6, text, size=28, color=BLACK, bold=True)

def add_sub(slide, left, top, text, width=10):
    add_text(slide, left, top, width, 0.4, text, size=13, color=GRAY)

def add_rect(slide, left, top, w, h, fill_color=HANA_LIGHT, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_hana_footer(slide, num, total):
    add_text(slide, 0.5, 7.0, 2, 0.3, f'{num} / {total}', size=9, color=LIGHT_GRAY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), W, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HANA
    bar.line.fill.background()

TOTAL = 21

# ===== SLIDE 1: HERO =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = HANA_DARK
bg_shape.line.fill.background()
add_text(s, 0.8, 1.5, 5, 0.3, '2026 하나 청년 금융인재 양성 프로젝트  |  Team UniHana', size=11, color=RGBColor(0xA8,0xE6,0xCF))
add_text(s, 0.8, 2.3, 10, 1.0, 'ESG TradeGuard', size=48, color=WHITE, bold=True)
add_text(s, 0.8, 3.5, 10, 0.5, 'AI 기반 수출환어음 ESG 컴플라이언스 심사 플랫폼', size=20, color=WHITE, bold=True)
add_text(s, 0.8, 4.3, 8, 1.0, '수출 서류 업로드부터 OCR 파싱, EU CBAM/EUDR 규제 매칭,\nCarbonCast 탄소비용 시뮬레이션, 위성 NDVI 환경 검증까지\n전 과정을 AI가 자동 처리합니다.', size=14, color=RGBColor(0xCC,0xCC,0xCC))
add_text(s, 0.8, 5.8, 5, 0.3, 'UniHana (유니하나)', size=12, color=RGBColor(0x99,0x99,0x99))
add_hana_footer(s, 1, TOTAL)

# ===== SLIDE 2: 목차 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG)
add_label(s, 1.5, 1.0, 'CONTENTS')
add_title(s, 1.5, 1.3, '목차')
toc_items = [
    ('01', '배경', '왜 이 문제를 다루는가'),
    ('02', '접근법', '거래 단위 ESG 분석'),
    ('03', '시스템 아키텍처', '기존 인프라 위에 ESG 레이어 추가'),
    ('04', '은행 심사 데모', 'OCR → 규제 매칭 → 위성 검증 → 리포트'),
    ('05', '기업 자가검진 데모', 'B2B SaaS 자가검진 + 비용 시뮬레이션'),
    ('06', '핵심 기능 + 시연 영상', '프로토타입 라이브 데모'),
    ('07', '경쟁 분석 + 수익 모델', '차별점과 비즈니스 구조'),
    ('08', '기대효과 + 로드맵', '구현 일정과 향후 확장'),
]
for i, (num, title, desc) in enumerate(toc_items):
    col = i % 2
    row = i // 2
    x = 1.5 + col * 5.2
    y = 2.3 + row * 1.15
    r = add_rect(s, x, y, 4.8, 0.95, WHITE, BORDER)
    add_text(s, x + 0.15, y + 0.15, 0.5, 0.4, num, size=13, color=HANA, bold=True)
    add_text(s, x + 0.7, y + 0.12, 3.5, 0.35, title, size=14, color=BLACK, bold=True)
    add_text(s, x + 0.7, y + 0.5, 3.5, 0.3, desc, size=10, color=GRAY)
add_hana_footer(s, 2, TOTAL)

# ===== SLIDE 3: 배경 Part 1 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG)
add_label(s, 0.8, 0.6, 'BACKGROUND')
add_title(s, 0.8, 0.9, '수출 거래의 비용 구조가 바뀌었습니다')
add_sub(s, 0.8, 1.5, 'EU CBAM 본격 시행 — 데이터 유무에 따라 같은 제품의 탄소비용이 두 배 차이')

# 두 카드
for ci, (title, items, bcolor) in enumerate([
    ('CBAM이 바꾼 것', ['기존: 관세나 인증처럼 통과 여부의 문제', 'CBAM: 데이터 유무에 따라 비용 자체가 달라지는 구조', '철강 기준 — 한국 실측 2.05 vs EU 기본값 3.90 tCO2/t', '데이터 유무만으로 탄소비용이 거의 두 배'], HANA),
    ('중소 수출기업의 현실', ['CBAM 영향권 수출기업 약 3,162개사', '데이터 인프라도 전담 인력도 없는 기업이 대다수', '직접 영향권 기업의 54.9%가 대응 계획조차 없음', 'EUDR 대응 완료 기업 12%, 중소기업은 5% 미만'], HANA_DARK),
]):
    x = 0.8 + ci * 5.9
    r = add_rect(s, x, 2.1, 5.5, 2.8, WHITE, BORDER)
    add_text(s, x + 0.3, 2.2, 4, 0.35, title, size=15, color=bcolor, bold=True)
    for j, item in enumerate(items):
        add_text(s, x + 0.3, 2.65 + j * 0.5, 4.8, 0.4, f'• {item}', size=11, color=GRAY)

# 핵심 질문 박스
r = add_rect(s, 0.8, 5.2, 11.7, 1.8, WHITE, HANA)
add_text(s, 1.0, 5.3, 11.3, 0.35, '핵심 질문', size=16, color=HANA, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 5.75, 11.3, 0.4, '왜 거래를 결정하는 단계에서는 이 정보가 반영되지 않는가?', size=15, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.25, 11.3, 0.6, '현재 수출환어음 매입 심사에서는 ESG 규제로 인한 추가 비용이나 리스크가 충분히 고려되지 않습니다.', size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_hana_footer(s, 3, TOTAL)

# ===== SLIDE 4: 배경 Part 2 - 규제 + 시장 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG)
add_label(s, 0.8, 0.6, 'BACKGROUND')
add_title(s, 0.8, 0.9, '주요 ESG 규제 현황 및 시장 규모')

regs = [
    ('EU CBAM', '철강, 알루미늄, 시멘트, 비료, 수소, 전력 6개 업종\n2026.01 본격 시행\n무상할당 폐지 2026년 2.5% → 2034년 100%'),
    ('EU EUDR', '소고기, 커피, 팜오일, 고무, 대두, 목재, 코코아 7개 품목\n산림 훼손 없는 공급망 증명 의무\n기업 자기 신고에 의존'),
    ('글로벌 확산', '영국 CBAM 2027년 시행 확정\n미국 청정경쟁법(CCA) 법안 발의\n일본 GX 국가전략 시행 중'),
]
for i, (t, d) in enumerate(regs):
    x = 0.8 + i * 4.0
    r = add_rect(s, x, 1.7, 3.7, 2.2, WHITE, BORDER)
    add_text(s, x + 0.2, 1.8, 3.2, 0.3, t, size=13, color=HANA, bold=True)
    add_text(s, x + 0.2, 2.2, 3.2, 1.5, d, size=10, color=GRAY)

# TAM/SAM/SOM
tams = [('TAM', '$68.3B', '한국 대EU 수출 전체'), ('SAM', '~800개사', '3,162개사 x 하나은행 25%'), ('SOM (1년차)', '50개사', 'B2B SaaS 월 30만원')]
for i, (label, val, desc) in enumerate(tams):
    x = 0.8 + i * 4.0
    r = add_rect(s, x, 4.3, 3.7, 2.5, WHITE, HANA)
    add_text(s, x + 0.2, 4.4, 3.2, 0.25, label, size=10, color=HANA, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 4.8, 3.2, 0.6, val, size=28, color=HANA, bold=True, align=PP_ALIGN.CENTER, font_name='JetBrains Mono')
    add_text(s, x + 0.2, 5.6, 3.2, 0.5, desc, size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_hana_footer(s, 4, TOTAL)

# ===== SLIDE 5: 접근법 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_label(s, 0.8, 0.6, 'OUR APPROACH')
add_title(s, 0.8, 0.9, '거래 단위 ESG 분석')
add_sub(s, 0.8, 1.5, '기존 ESG 분석은 기업 단위 등급 중심. 비용은 거래에서 발생합니다.')

r = add_rect(s, 0.8, 2.1, 11.7, 1.2, HANA_LIGHT, HANA)
add_text(s, 1.1, 2.2, 11, 0.3, '문제 재정의', size=14, color=HANA_DARK, bold=True)
add_text(s, 1.1, 2.6, 11, 0.6, '같은 기업이라도 수출 품목과 목적지에 따라 CBAM 비용과 규제 리스크는 달라집니다.\n분석 단위를 \'수출 거래 한 건\'으로 전환하는 것이 핵심입니다.', size=12, color=BLACK)

approaches = [
    ('방안 1. 규제 정보 추출', 'AI-OCR로 5종 서류에서 HS코드, 수출량,\n원산지 등 15개 핵심 필드를 추출하고\nCBAM/EUDR 적용 여부를 자동 판정'),
    ('방안 2. 탄소 데이터 신뢰도', 'CarbonCast 엔진에 z-score 기반\n4단계 감점 로직 적용, 배출 데이터\n품질을 0~100점으로 평가'),
    ('방안 3. 환경 리스크 보완', 'Sentinel-2 위성 NDVI 데이터로\n원산지 변화를 1차 확인,\nFSC/PEFC 인증 API로 교차 검증'),
]
for i, (t, d) in enumerate(approaches):
    x = 0.8 + i * 4.0
    r = add_rect(s, x, 3.6, 3.7, 2.5, WHITE, BORDER)
    add_text(s, x + 0.2, 3.7, 3.2, 0.3, t, size=12, color=HANA, bold=True)
    add_text(s, x + 0.2, 4.1, 3.2, 1.8, d, size=10, color=GRAY)

r1 = add_rect(s, 0.8, 6.4, 5.7, 0.5, BG, BORDER)
add_text(s, 1.0, 6.42, 5.3, 0.4, '은행 심사용 — 무역금융 담당자가 수출환어음 매입 심사 시 사용', size=10, color=GRAY)
r2 = add_rect(s, 6.8, 6.4, 5.7, 0.5, BG, BORDER)
add_text(s, 7.0, 6.42, 5.3, 0.4, '기업 자가진단용 — 수출기업이 직접 ESG 규제 리스크를 사전 점검', size=10, color=GRAY)
add_hana_footer(s, 5, TOTAL)

# ===== SLIDE 6: 아키텍처 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_label(s, 0.8, 0.6, 'ARCHITECTURE')
add_title(s, 0.8, 0.9, '시스템 아키텍처')
add_sub(s, 0.8, 1.5, '기존 인프라 위에 ESG 레이어만 추가')

steps = ['입력 레이어\n수출 인보이스\n원산지증명서\nB/L · 신용장', 'READIT OCR\nGS인증 1등급\n다국어 서류 인식\nHS코드 자동 추출', 'AI 분석 엔진\nTransformer NER\nLLM 규제 매칭\nESG 규제 DB 조회', '위성 검증\nSentinel-2 NDVI\nFSC/PEFC API\n하이브리드 검증', '리스크 보고서\n3단계 등급\n조치사항 체크리스트\n담당자 최종 판단']
for i, step in enumerate(steps):
    x = 0.5 + i * 2.5
    r = add_rect(s, x, 2.2, 2.1, 2.5, WHITE, BORDER)
    lines = step.split('\n')
    add_text(s, x + 0.15, 2.3, 1.8, 0.3, lines[0], size=11, color=HANA, bold=True, align=PP_ALIGN.CENTER)
    for j, line in enumerate(lines[1:]):
        add_text(s, x + 0.15, 2.7 + j * 0.4, 1.8, 0.3, line, size=9, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s, x + 2.15, 3.2, 0.3, 0.3, '→', size=16, color=HANA, bold=True)

infra = [('ESG 규제 DB', 'EU CBAM/EUDR 규정, HS코드-규제 매핑'), ('하나은행 시스템 연동', '비대면 AI 수출환어음매입 심사 시스템과 API 연동'), ('CarbonCast 엔진', 'z-score 4단계 감점, NGFS 시나리오 시뮬레이션')]
for i, (t, d) in enumerate(infra):
    x = 0.8 + i * 4.0
    r = add_rect(s, x, 5.2, 3.7, 1.2, HANA_LIGHT)
    add_text(s, x + 0.2, 5.3, 3.2, 0.3, t, size=12, color=HANA_DARK, bold=True)
    add_text(s, x + 0.2, 5.7, 3.2, 0.5, d, size=10, color=GRAY)
add_hana_footer(s, 6, TOTAL)

# ===== SLIDE 7~11: 은행 심사 데모 (간략 텍스트 슬라이드) =====
bank_slides = [
    ('서류 업로드', '수출 관련 서류를 드래그앤드롭으로 업로드', [
        'Commercial_Invoice.pdf — 2.4 MB ✓',
        'Certificate_of_Origin_Final.pdf — 1.8 MB ✓',
        'Bill_of_Lading.pdf — 1.1 MB ✓',
        'Packing_List.pdf — 890 KB ✓',
        'Carbon_Emission_Report.pdf — 2.1 MB ✓',
    ], '은행 내부용 — 무역금융 심사'),
    ('OCR 자동 파싱', 'AI-OCR로 5종 서류에서 15개 핵심 필드 자동 추출', [
        '수출자: Northwind Corp., San Francisco, CA',
        '수입자: Contoso GmbH, Berlin, Germany (EU)',
        'HS코드: 8471.30 — Portable Computing Devices (외 5건)',
        '거래 금액: USD 379,320.00 (CIF Hamburg)',
        '원산지: United States of America',
        '탄소배출량: 1.42 tonnes CO₂e (GLEC Framework v3.0)',
    ], None),
    ('ESG 규제 자동 매칭', 'HS코드 기반 CBAM/EUDR 해당 여부 판단', [
        'EU CBAM — CLEAR: HS 8471 (전자기기)은 CBAM 대상 아님',
        'EU EUDR — CLEAR: IT 전자기기는 EUDR 대상 품목 아님',
        'EU RoHS/WEEE — REVIEW: RoHS 적합성 선언서 미첨부',
        '수출통제/제재 — CLEAR: EAR99 일반 상용 등급',
    ], None),
    ('위성 NDVI 환경 검증', 'Sentinel-2 위성 1차 스크리닝 → FSC/PEFC API 최종 검증', [
        '위치: Oakland, CA — Port of Oakland 인근',
        'NDVI 지수: 0.43 (항만 산업지구 정상 범위)',
        '12개월 추이: 0.41~0.45 안정적',
        '판정: 산림벌채 관련 리스크 해당 없음',
    ], None),
    ('최종 리스크 보고서', 'Northwind Corp. → Contoso GmbH · USD 379,320', [
        '종합 점수: 72/100 (조건부 승인)',
        'EU CBAM — CLEAR (현행 기준 리스크 없음)',
        'EU RoHS/WEEE — REVIEW (RoHS 적합성 선언서 미첨부)',
        '운송 탄소배출 — LOW (1.42 tCO₂e)',
        '조치: RoHS 선언서 제출 요청, WEEE 등록번호 확인',
    ], None),
]

for idx, (title, sub, items, extra) in enumerate(bank_slides):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    if extra:
        r = add_rect(s, 0.8, 0.5, 11.7, 0.6, HANA_LIGHT)
        add_text(s, 1.0, 0.55, 10, 0.35, extra, size=12, color=HANA_DARK, bold=True)
        y_start = 1.3
    else:
        y_start = 0.6
    add_text(s, 0.8, y_start, 2, 0.35, f'STEP {idx+1}', size=10, color=WHITE, bold=True)
    # step badge
    badge = add_rect(s, 0.8, y_start, 0.8, 0.3, HANA)
    tf = badge.text_frame
    tf.paragraphs[0].text = f'STEP {idx+1}'
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_title(s, 0.8, y_start + 0.5, title)
    add_sub(s, 0.8, y_start + 1.1, sub)

    # 스크린 프레임
    frame = add_rect(s, 0.8, y_start + 1.7, 11.7, 4.5 - (0.7 if extra else 0), WHITE, BORDER)
    for j, item in enumerate(items):
        check = '✓' if '✓' in item or 'CLEAR' in item else '⚠' if 'REVIEW' in item else '•'
        c = HANA if check == '✓' else HANA_DARK if check == '⚠' else GRAY
        add_text(s, 1.2, y_start + 2.0 + j * 0.55, 10.5, 0.4, item, size=12, color=c)
    add_hana_footer(s, 7 + idx, TOTAL)

# ===== SLIDE 12~15: 기업 자가검진 데모 =====
biz_slides = [
    ('서류 업로드', '(주) 그린텍 · HS 7604.29 (알루미늄 프로파일) · 프랑스 (EU)', [
        'Invoice_그린텍_2026Q1.pdf — 1.8 MB ✓',
        'Carbon_Report_그린텍_2025.pdf — 2.1 MB ✓',
        'Certificate_of_Origin_KR.pdf — 640 KB ✓',
    ], '기업 자가검진용 — B2B SaaS'),
    ('자가검진 결과', 'HS 7604.29 (알루미늄 프로파일) · 프랑스 (EU)', [
        '종합 점수: 81/100 (양호)',
        'EU CBAM — REVIEW: 알루미늄은 CBAM 대상, 벤치마크 근접',
        'EU EUDR — CLEAR: 알루미늄은 EUDR 대상 아님',
        '탄소 집약도 — PASS: 1.48 tCO2/t (벤치마크 1.52 이내)',
        '위성 NDVI — CLEAR: NDVI 0.71 안정적',
    ], None),
    ('규제 가이드', 'HS 7604.29 (알루미늄 프로파일) · 프랑스 (EU)', [
        'EU CBAM — 대상 품목: Regulation (EU) 2023/956, Annex I',
        '의무: CBAM 인증서 구매, 내재 탄소배출량 신고, 분기별 보고서',
        '벤치마크: 알루미늄 1.52 tCO2/t (EU ETS 무상할당 기준)',
        'EU EUDR — 비대상: 알루미늄은 EUDR 7개 품목 해당 없음',
        '타임라인: 서류 준비(Now) → CBAM 보고서(Q2) → 인증서 구매 → 수출 통관',
    ], None),
    ('비용 시뮬레이션', 'HS 7604.29 (알루미늄 프로파일) · 프랑스 (EU)', [
        '수출 수량: 120 MT | 실제 배출량: 1.48 tCO2/t | 벤치마크: 1.52',
        '현재 (2026): 예상 추가 비용 EUR 0',
        '보수적 (2027): 벤치마크 1.45 → EUR 252',
        '공격적 (2028): 벤치마크 1.35 → EUR 1,248',
        '권고: 탄소 저감 계획 수립, 에너지 효율 개선 투자 검토',
    ], None),
]

for idx, (title, sub, items, extra) in enumerate(biz_slides):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    if extra:
        r = add_rect(s, 0.8, 0.5, 11.7, 0.6, HANA_LIGHT)
        add_text(s, 1.0, 0.55, 10, 0.35, extra, size=12, color=HANA_DARK, bold=True)
        y_start = 1.3
    else:
        y_start = 0.6
    badge = add_rect(s, 0.8, y_start, 0.8, 0.3, HANA)
    tf = badge.text_frame
    tf.paragraphs[0].text = f'STEP {idx+1}'
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_title(s, 0.8, y_start + 0.5, title)
    add_sub(s, 0.8, y_start + 1.1, sub)

    frame = add_rect(s, 0.8, y_start + 1.7, 11.7, 4.5 - (0.7 if extra else 0), WHITE, BORDER)
    for j, item in enumerate(items):
        c = HANA if '✓' in item or 'CLEAR' in item or 'PASS' in item or 'EUR 0' in item else HANA_DARK if 'REVIEW' in item or '권고' in item else GRAY
        add_text(s, 1.2, y_start + 2.0 + j * 0.55, 10.5, 0.4, item, size=12, color=c)
    add_hana_footer(s, 12 + idx, TOTAL)

# ===== SLIDE 16: 핵심 기능 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG)
add_label(s, 0.8, 0.6, 'FEATURES')
add_title(s, 0.8, 0.9, '핵심 기능')
add_sub(s, 0.8, 1.5, '기존 수출환어음 심사 시스템 위에 ESG 레이어를 추가합니다')

features = [
    ('READIT OCR 파싱', 'GS인증 1등급 AI-OCR로 수출 서류를\n자동 인식. 98%+ 정확도로 추출'),
    ('LLM 규제 매칭', 'HS코드 기반 EU CBAM, EUDR 등\nESG 규제 해당 여부를 자동 판단'),
    ('위성 NDVI 검증', 'Sentinel-2 위성 데이터로 원산지\n환경 영향을 1차 스크리닝'),
    ('리스크 등급 시각화', '3단계 리스크 등급으로\n직관적 판단 지원'),
    ('Transformer NER', '무역 서류 내 엔티티를\nTransformer 기반 NER로 관계 추출'),
    ('규제 자동 업데이트', 'EU 관보 웹크롤링으로\n규제 변경 자동 감지'),
]
for i, (t, d) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 2.1 + row * 2.5
    r = add_rect(s, x, y, 3.7, 2.1, WHITE, BORDER)
    add_text(s, x + 0.2, y + 0.2, 3.2, 0.3, t, size=13, color=BLACK, bold=True)
    add_text(s, x + 0.2, y + 0.6, 3.2, 1.2, d, size=10, color=GRAY)
add_hana_footer(s, 16, TOTAL)

# ===== SLIDE 17: 시연 영상 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_label(s, 1.5, 1.5, 'LIVE DEMO')
add_title(s, 1.5, 1.8, '프로토타입 시연 영상', width=10)
add_sub(s, 1.5, 2.4, '실제 작동 프로토타입 시연')
# 비디오 플레이스홀더
r = add_rect(s, 1.5, 3.2, 10.3, 3.5, RGBColor(0x1A,0x1A,0x1A), BORDER)
add_text(s, 1.5, 4.5, 10.3, 0.5, '▶  demo-video.mp4', size=18, color=RGBColor(0x66,0x66,0x66), align=PP_ALIGN.CENTER)
add_hana_footer(s, 17, TOTAL)

# ===== SLIDE 18: 경쟁 분석 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_label(s, 0.8, 0.6, 'COMPETITIVE EDGE')
add_title(s, 0.8, 0.9, '경쟁 분석')
add_sub(s, 0.8, 1.5, '기존 ESG 솔루션과의 차별점')

# 테이블 대신 카드
headers = ['비교 항목', '기존 ESG 솔루션', '컨설팅사', 'ESG TradeGuard']
rows = [
    ['영역', 'ESG 공시/평가', 'CBAM 컨설팅', '무역금융 실시간 ESG 심사'],
    ['타이밍', '사후 평가', '프로젝트 단위', '수출 거래 시점 실시간'],
    ['비용', '연간 구독료', '건당 수천만원', '수출환어음 수수료에 포함'],
    ['접근성', '대기업 위주', '대기업 위주', '모든 수출 중소기업'],
    ['금융 연계', '없음', '없음', '수출환어음 매입과 동시 처리'],
]
# Header row
for j, h in enumerate(headers):
    x = 0.8 + j * 3.0
    c = HANA if j == 3 else HANA_DARK
    r = add_rect(s, x, 2.1, 2.8, 0.45, c)
    add_text(s, x + 0.1, 2.13, 2.6, 0.35, h, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
# Data rows
for i, row in enumerate(rows):
    for j, cell in enumerate(row):
        x = 0.8 + j * 3.0
        y = 2.6 + i * 0.55
        bg_c = HANA_LIGHT if j == 3 else WHITE
        r = add_rect(s, x, y, 2.8, 0.5, bg_c, BORDER)
        c = HANA_DARK if j == 3 else (BLACK if j == 0 else GRAY)
        b = j == 0 or j == 3
        add_text(s, x + 0.1, y + 0.05, 2.6, 0.35, cell, size=10, color=c, bold=b, align=PP_ALIGN.CENTER)

# 3 강점 카드
strengths = [
    ('거래 시점 실시간 검증', '수출 서류 제출 → 즉시 ESG 규제 체크'),
    ('금융과 ESG의 자연스러운 융합', '기존 수출환어음 심사 프로세스에 ESG 레이어 추가'),
    ('중소기업 접근성', '컨설팅 비용 없이 은행 거래만으로 ESG 컴플라이언스 확인'),
]
for i, (t, d) in enumerate(strengths):
    x = 0.8 + i * 4.0
    r = add_rect(s, x, 5.6, 3.7, 1.2, HANA_LIGHT)
    add_text(s, x + 0.2, 5.7, 3.2, 0.3, t, size=11, color=HANA_DARK, bold=True)
    add_text(s, x + 0.2, 6.1, 3.2, 0.5, d, size=9, color=GRAY)
add_hana_footer(s, 18, TOTAL)

# ===== SLIDE 19: 수익 모델 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG)
add_label(s, 0.8, 0.6, 'REVENUE MODEL')
add_title(s, 0.8, 0.9, '수익 모델')

metrics = [('수수료 수익', '수출환어음 매입\n수수료 증대'), ('B2B SaaS', '기업 자가검진\n월 30만원'), ('고객 록인', '수출기업 →\n외환 거래 수익'), ('ESG 우대', 'ESG 우수기업\n우대금리 대출')]
for i, (t, d) in enumerate(metrics):
    x = 0.8 + i * 3.0
    r = add_rect(s, x, 1.7, 2.8, 1.3, WHITE, BORDER)
    add_text(s, x + 0.1, 1.8, 2.6, 0.35, t, size=14, color=HANA, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 2.3, 2.6, 0.5, d, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# 직접/간접 수익
for ci, (title, items) in enumerate([
    ('직접 수익', ['수출환어음 매입 수수료: ESG 체크 포함 시 프리미엄', 'B2B SaaS 구독: 기업 자가검진 서비스 월 30만원', 'ESG 컴플라이언스 리포트: 기업별 맞춤 보고서 유료 제공', 'ESG 우수 기업 우대 금리 → 대출 취급 증가']),
    ('간접 수익', ['수출기업 고객 유입: 경쟁사 대비 차별화', '외환 거래 수익: 수출기업 록인 → 환전/파생상품 증가', '무역보험 연계: ESG 리스크 데이터 기반 상품 개발', '글로벌 확장: 영국 CBAM 2027, 미국 CCA, 일본 GX']),
]):
    x = 0.8 + ci * 6.0
    r = add_rect(s, x, 3.3, 5.7, 3.5, WHITE, BORDER)
    add_text(s, x + 0.3, 3.4, 5, 0.35, title, size=14, color=HANA_DARK, bold=True)
    for j, item in enumerate(items):
        add_text(s, x + 0.3, 3.9 + j * 0.6, 5, 0.5, f'• {item}', size=10, color=GRAY)
add_hana_footer(s, 19, TOTAL)

# ===== SLIDE 20: 기대효과 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG)
add_label(s, 0.8, 0.6, 'EXPECTED IMPACT')
add_title(s, 0.8, 0.9, '기대효과')
add_sub(s, 0.8, 1.5, '거래별 비용과 리스크를 심사 단계에서 함께 고려')

for ci, (title, items) in enumerate([
    ('은행 심사 측면', ['거래별 예상 비용이 함께 제시 → 금리와 한도 설정에 반영', '추가 비용이 큰 거래 → 금리 보수적 적용 또는 한도 축소', '비용 부담이 낮은 거래 → 금리 인하 또는 한도 확대', 'ESG 미준수 기업 수출환어음 매입 시 은행 리스크 사전 차단']),
    ('수출기업 측면', ['거래 이전 단계에서 예상 비용과 리스크를 사전 확인', '거래 진행/조건 조정/추가 데이터 준비 결정 가능', '준비 수준에 따라 비용 차이를 수치로 확인', '대응 방향을 구체적으로 설정할 수 있음']),
]):
    x = 0.8 + ci * 6.0
    r = add_rect(s, x, 2.1, 5.7, 3.0, WHITE, BORDER)
    add_text(s, x + 0.3, 2.2, 5, 0.35, title, size=14, color=HANA_DARK, bold=True)
    for j, item in enumerate(items):
        add_text(s, x + 0.3, 2.7 + j * 0.55, 5, 0.45, f'• {item}', size=10, color=GRAY)

extras = [('데이터 자산화', '거래 단위 데이터를 구독형 리스크 리포트로 확장'), ('금융상품 연계', 'ESG 연계 금융상품, 맞춤형 컨설팅 서비스'), ('즉시 도입 가능', '기존 프로세스에 모듈 형태로 추가')]
for i, (t, d) in enumerate(extras):
    x = 0.8 + i * 4.0
    r = add_rect(s, x, 5.5, 3.7, 1.2, WHITE, BORDER)
    add_text(s, x + 0.2, 5.6, 3.2, 0.3, t, size=12, color=HANA_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 6.0, 3.2, 0.5, d, size=9, color=GRAY, align=PP_ALIGN.CENTER)
add_hana_footer(s, 20, TOTAL)

# ===== SLIDE 21: 로드맵 + 마무리 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_label(s, 0.8, 0.6, 'ROADMAP')
add_title(s, 0.8, 0.9, '구현 로드맵')
add_sub(s, 0.8, 1.5, '6개월 내 프로토타입, 12개월 내 상용화')

timeline = [
    ('Month 1~2', '규제 DB 구축', 'EU CBAM/EUDR 규정 분석, HS코드-규제 매핑'),
    ('Month 2~3', 'OCR + LLM 파이프라인', '수출 서류 OCR 파싱, LLM 프롬프트 엔지니어링'),
    ('Month 3~4', '프로토타입 통합', '웹 UI 개발, 전체 파이프라인 통합'),
    ('Month 4~5', '검증 및 고도화', '실제 수출 서류 기반 테스트, 정확도 개선'),
    ('Month 5~6', '최종 데모', '라이브 데모 준비, B2B SaaS 모듈 완성'),
]
for i, (time, title, desc) in enumerate(timeline):
    y = 2.1 + i * 0.75
    # 타임라인 도트
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(y + 0.1), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = HANA
    dot.line.fill.background()
    if i < len(timeline) - 1:
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.265), Inches(y + 0.25), Inches(0.02), Inches(0.55))
        line.fill.solid()
        line.fill.fore_color.rgb = HANA
        line.line.fill.background()
    add_text(s, 1.6, y, 2, 0.3, time, size=11, color=HANA, bold=True)
    add_text(s, 3.5, y, 3, 0.3, title, size=12, color=BLACK, bold=True)
    add_text(s, 3.5, y + 0.3, 4, 0.3, desc, size=9, color=GRAY)

# 기술 스택
techs = [('OCR', 'READIT OCR'), ('LLM', 'GPT-4 / Claude'), ('NER', 'Transformer'), ('위성', 'Sentinel-2'), ('인증', 'FSC/PEFC API'), ('프론트', 'React + TS'), ('백엔드', 'Python FastAPI')]
for i, (k, v) in enumerate(techs):
    y = 2.1 + i * 0.55
    add_text(s, 8.5, y, 2, 0.25, k, size=9, color=GRAY)
    add_text(s, 10.0, y, 2.5, 0.25, v, size=10, color=BLACK, bold=True)

# 마무리 배너
r = add_rect(s, 0.8, 6.0, 11.7, 0.9, HANA_DARK)
add_text(s, 1.0, 6.1, 11.3, 0.35, '2026년 CBAM 본격 시행, 수출 심사 시점에 ESG 규제를 자동 검증합니다', size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.5, 11.3, 0.25, 'ESG TradeGuard — Team UniHana', size=10, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.CENTER)
add_hana_footer(s, 21, TOTAL)

# ===== SAVE =====
prs.save('ESG_TradeGuard_Presentation.pptx')
print('Done: ESG_TradeGuard_Presentation.pptx')
